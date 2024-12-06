from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def inspect_shape(shape, indent=0):
    """Recursively inspect a shape and print its text content."""
    if shape.has_text_frame:
        print(" " * indent + f"Shape: {shape.shape_type} - Text: '{shape.text}'")
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                print(" " * (indent + 2) + f"Run: '{run.text}'")
    elif shape.has_table:
        print(" " * indent + f"Table: {shape.table}")
        for row in shape.table.rows:
            for cell in row.cells:
                print(" " * (indent + 2) + f"Cell: '{cell.text}'")
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # Check for grouped shapes
        print(" " * indent + f"Grouped Shape: {shape.shape_type}")
        for s in shape.shapes:
            inspect_shape(s, indent + 2)  # Recursively inspect inner shapes
    else:
        print(" " * indent + f"Other Shape Type: {shape.shape_type}")

def inspect_presentation(file_path):
    presentation = Presentation(file_path)
    for slide_num, slide in enumerate(presentation.slides, start=1):
        print(f"Slide {slide_num}:")
        for shape in slide.shapes:
            inspect_shape(shape)  # Inspect each shape

if __name__ == "__main__":
    inspect_presentation("translated_es (1).pptx") 