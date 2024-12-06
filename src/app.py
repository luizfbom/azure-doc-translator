from pptx import Presentation
from services.azure_service import translate_pptx
import os
import streamlit as st
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_from_presentation(presentation):
    text_list = []
    st.write("Extracting text from presentation...")

    def extract_text_from_shape(shape):
        """Recursively extract text from a shape or group of shapes."""
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():  # Only add non-empty text
                        text_list.append(run.text.strip())  # Ensure we strip whitespace
        elif shape.has_table:  # Check for text in tables
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_list.append(cell.text.strip())
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # Check for grouped shapes
            for s in shape.shapes:
                extract_text_from_shape(s)  # Recursively extract text from grouped shapes

    for slide_num, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            extract_text_from_shape(shape)  # Extract text from each shape

    # Remove duplicates while preserving order
    unique_text_list = list(dict.fromkeys(text_list))
    return unique_text_list

def translate_presentation(presentation, translated_phrase_list):
    # Translate slides
    for slide in presentation.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, translated_phrase_list)

def replace_text_in_shape(shape, translated_phrase_list):
    """Recursively replace text in a shape or group of shapes."""
    if shape.has_text_frame:
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            # Combine runs to check for full text matches
            full_text = ''.join(run.text for run in paragraph.runs).strip()
            # Check if the full text is in the translation dictionary
            if full_text in translated_phrase_list:
                translated_text = translated_phrase_list[full_text]
                # Replace all runs with the translated text
                for run in paragraph.runs:
                    run.text = translated_text
            else:
                # Check each run individually for partial matches
                for run in paragraph.runs:
                    original_text = run.text.strip()
                    if original_text in translated_phrase_list:
                        run.text = translated_phrase_list[original_text]
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # Check for grouped shapes
        for s in shape.shapes:
            replace_text_in_shape(s, translated_phrase_list)  # Recursively replace text in grouped shapes

def main():
    st.title("PowerPoint Translator")
    st.write("Upload a PowerPoint file and select a target language for translation")

    # File uploader
    uploaded_file = st.file_uploader("Choose a PPTX file", type="pptx")
    
    # Dropdown for language selection
    language_options = {
        "English": "en",
        "Spanish": "es",
        "French": "fr",
        "German": "de",
        "Italian": "it",
        "Portuguese": "pt",
        "Chinese": "zh",
        "Japanese": "ja",
        "Russian": "ru",
        "Arabic": "ar"
    }
    
    target_language = st.selectbox("Select target language", list(language_options.keys()))
    
    if uploaded_file and target_language and st.button("Translate"):
        temp_file = "temp.pptx"
        
        # Get the original filename without extension
        original_filename = os.path.splitext(uploaded_file.name)[0]
        translated_file_path = f"{original_filename}_translated_{language_options[target_language]}.pptx"
        
        try:
            # Save uploaded file temporarily
            with open(temp_file, "wb") as f:
                f.write(uploaded_file.getvalue())
            
            # Process the presentation
            presentation = Presentation(temp_file)
            phrase_list = extract_text_from_presentation(presentation)
            
            # Show translation progress
            with st.spinner('Translating...'):
                translated_phrase_list = translate_pptx(phrase_list, language_options[target_language])
            
            if translated_phrase_list:
                translate_presentation(presentation, translated_phrase_list)
                presentation.save(translated_file_path)
                
                # Provide download link
                with open(translated_file_path, "rb") as file:
                    st.download_button(
                        label="Download translated presentation",
                        data=file,
                        file_name=translated_file_path,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")
        
        finally:
            # Clean up all temporary files
            for file_path in [temp_file, translated_file_path]:
                try:
                    if os.path.exists(file_path):
                        os.remove(file_path)
                except Exception as e:
                    print(f"Error deleting {file_path}: {e}")

if __name__ == "__main__":
    main()